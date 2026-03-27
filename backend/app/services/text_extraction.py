"""Text extraction from documents and token-aware chunking."""
from __future__ import annotations

import io
import re

import structlog

logger = structlog.get_logger()

SUPPORTED_TEXT_TYPES = {
    "text/plain",
    "text/markdown",
    "text/csv",
    "text/html",
    "application/json",
    "application/xml",
    "text/xml",
}

PDF_TYPES = {"application/pdf"}

DOCX_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}

XLSX_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}

PPTX_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def extract_text_from_file(
    file_bytes: bytes,
    content_type: str,
    filename: str,
) -> str:
    """Extract plain text from a document file.

    Returns empty string for unsupported types (embedding is skipped).
    """
    ct = content_type.lower().strip()

    try:
        if ct in PDF_TYPES or filename.lower().endswith(".pdf"):
            return _extract_pdf(file_bytes)

        if ct in DOCX_TYPES or filename.lower().endswith(".docx"):
            return _extract_docx(file_bytes)

        if ct in XLSX_TYPES or filename.lower().endswith(".xlsx"):
            return _extract_xlsx(file_bytes)

        if ct in PPTX_TYPES or filename.lower().endswith(".pptx"):
            return _extract_pptx(file_bytes)

        if ct in SUPPORTED_TEXT_TYPES or ct.startswith("text/"):
            return file_bytes.decode("utf-8", errors="replace")

        # Code files by extension
        code_exts = {".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".go", ".rs",
                     ".rb", ".sh", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".sql"}
        ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext in code_exts:
            return file_bytes.decode("utf-8", errors="replace")

    except Exception:
        logger.warning("text_extraction_failed", filename=filename, content_type=content_type)
        return ""

    return ""


def _strip_markdown_formatting(text: str) -> str:
    """Lightly clean markdown for better embedding quality."""
    text = re.sub(r"```[\s\S]*?```", lambda m: m.group(0), text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(
    text: str,
    max_tokens: int = 800,
    overlap_tokens: int = 100,
) -> list[str]:
    """Split text into overlapping chunks sized by token count.

    Uses tiktoken cl100k_base encoding (GPT-4 / text-embedding-ada-002).
    """
    import tiktoken

    text = _strip_markdown_formatting(text)
    if not text:
        return []

    enc = tiktoken.get_encoding("cl100k_base")
    tokens = enc.encode(text)

    if len(tokens) <= max_tokens:
        return [text]

    chunks: list[str] = []
    start = 0
    step = max_tokens - overlap_tokens

    while start < len(tokens):
        end = min(start + max_tokens, len(tokens))
        chunk_tokens = tokens[start:end]
        chunk_text_str = enc.decode(chunk_tokens).strip()
        if chunk_text_str:
            chunks.append(chunk_text_str)
        if end >= len(tokens):
            break
        start += step

    return chunks
